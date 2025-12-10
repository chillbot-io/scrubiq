"""PowerPoint presentation (.pptx) extraction."""

from pathlib import Path

from .base import Extractor, ExtractionError

try:
    from pptx import Presentation

    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


class PptxExtractor(Extractor):
    """Extract text from PowerPoint presentations."""

    @property
    def extensions(self) -> list[str]:
        return [".pptx"]

    def extract(self, path: Path) -> str:
        """Extract text from slides, notes, and tables."""
        if not HAS_PPTX:
            raise ExtractionError("python-pptx not installed. Run: pip install python-pptx")

        try:
            prs = Presentation(path)
            text_parts = []

            for i, slide in enumerate(prs.slides):
                text_parts.append(f"[Slide {i + 1}]")

                # Extract from shapes (text boxes, titles, etc.)
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                text_parts.append(text)

                    # Extract from tables
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_text.append(cell.text.strip())
                            if row_text:
                                text_parts.append(" | ".join(row_text))

                # Extract from notes
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame
                    if notes and notes.text.strip():
                        text_parts.append("[Notes]")
                        text_parts.append(notes.text)

            return "\n".join(text_parts)

        except Exception as e:
            raise ExtractionError(f"Failed to extract from {path}: {e}")
